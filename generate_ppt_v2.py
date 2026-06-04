#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
水库安全生产汇报 PPT - 根据原文制作
政府汇报风格 - 简约正式
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色
COLOR_PRIMARY = RGBColor(25, 55, 109)
COLOR_SECONDARY = RGBColor(68, 114, 196)
COLOR_TEXT = RGBColor(0, 0, 0)
COLOR_TEXT_LIGHT = RGBColor(80, 80, 80)

def set_font(paragraph, size=18, bold=False, color=COLOR_TEXT):
    for run in paragraph.runs:
        run.font.name = '微软雅黑'
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

def add_title_bar(slide, title):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), prs.slide_width - Inches(1.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_font(p, size=24, bold=True, color=RGBColor(255, 255, 255))

def add_content(slide, items, top=Inches(1.2), height=None, font_size=16):
    if height is None:
        height = prs.slide_height - top - Inches(0.5)
    box = slide.shapes.add_textbox(Inches(0.6), top, prs.slide_width - Inches(1.2), height)
    tf = box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.space_after = Pt(8)
        set_font(p, size=font_size, color=COLOR_TEXT)
    return tf

# 第 1 页：封面
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.5))
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_PRIMARY
bg.line.fill.background()

title = slide.shapes.add_textbox(Inches(0.8), Inches(2), prs.slide_width - Inches(1.6), Inches(1.5))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "守牢安全底线 护航水务发展"
p.alignment = PP_ALIGN.CENTER
set_font(p, size=40, bold=True, color=COLOR_PRIMARY)

sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), prs.slide_width - Inches(1.6), Inches(0.8))
tf = sub.text_frame
p = tf.paragraphs[0]
p.text = "承德市水务局安全生产工作汇报"
p.alignment = PP_ALIGN.CENTER
set_font(p, size=20, color=COLOR_TEXT_LIGHT)

# 第 2 页：目录
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "汇报目录")
add_content(slide, [
    "一、切实加强水务安全生产工作的认识",
    "",
    "二、水务安全生产现状",
    "",
    "三、面临的形势",
    "",
    "四、下一步重点工作",
    "",
    "五、常见安全隐患"
], font_size=22)

# 第 3 页：一（一）书记重要论述
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "一、切实加强水务安全生产工作的认识")
add_content(slide, [
    "（一）书记关于安全生产重要论述核心要义",
    "",
    "书记始终将安全生产作为极重的民生大事，纳入"五位一体"总体布局和"四个全面"战略布局统筹推进",
    "",
    "核心要牢牢把握六个方面：",
    "1. 树牢红线意识，坚守"发展决不能以牺牲人的生命为代价"这条不可逾越的底线",
    "2. 健全责任体系，严格落实"党政同责、一岗双责、齐抓共管、失职追责"要求",
    "3. 压实主体责任，督促生产经营单位履行安全生产法定职责，做到安全投入、管理、培训、防护、应急"五到位"",
    "4. 创新监管方式，加快推进安全监管改革，提升监管精准度和效能",
    "5. 构建长效机制，实现风险常态化管控、隐患动态化治理",
    "6. 强化责任担当，以"时时放心不下"的责任感抓细抓实安全生产各项工作"
], font_size=15)

# 第 4 页：一（二）法律法规
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "（二）核心法律法规与政策文件")
add_content(slide, [
    "水务安全生产工作必须在法治框架下推进：",
    "",
    "通用层面：",
    "• 《安全生产法》明确了企业主体责任、政府监管职责、隐患排查治理和应急处置等基本原则",
    "",
    "水利领域：",
    "• 《水法》《防洪法》等法规对水利工程设计、施工、运行、维护全生命周期安全管理作出规范",
    "• 强调"事前防控、过程控制、事后处置"的闭环管理",
    "",
    "省市制度：",
    "• 《安全生产末端落实工作机制》",
    "• 《安全生产违法行为行政处罚规定》等制度，要求更严、标准更高、落点更实"
], font_size=16)

# 第 5 页：一（三）工作目标
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "（三）水务安全生产的目标")
add_content(slide, [
    "省水利厅总体目标：",
    "坚决杜绝重特大事故，遏制较大事故，减少一般事故，守牢安全底线",
    "",
    "我市水务安全生产工作更要聚焦实际、靶向发力：",
    "",
    "牢牢守住"人员不伤亡、财产少损失"的基本底线",
    "",
    "确保全市水利工程建设有序推进",
    "",
    "确保各类水利设施安全稳定运行",
    "",
    "为全市经济社会发展提供坚实的水安全保障"
], font_size=18)

# 第 6 页：二（一）事故统计
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "二、水务安全生产现状 - 事故统计")
add_content(slide, [
    "（一）事故统计数据",
    "",
    "从全国范围看，2021-2023 年水利行业共发生生产安全事故 136 起，造成 152 人死亡、28 人受伤",
    "",
    "其中：",
    "• 水利建设项目事故占比 68%，主要集中在高处坠落、物体打击、溺水、机械伤害等类型",
    "",
    "• 水库及水电站运行事故占比 22%，以设施故障引发的溃坝、触电等事故为主",
    "",
    "这些数据充分说明，水利行业安全生产风险依然突出，必须时刻绷紧安全这根弦"
], font_size=17)

# 第 7 页：二（二）安全生产特点
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "（二）安全生产特点")
add_content(slide, [
    "基于承德地理环境和水务工作实际，主要呈现四大特点：",
    "",
    "一是风险点多面广。我市地处"四河之源、两库上游"，境内有滦河、潮白河等四大水系，1500 条河流、众多水库和水电站",
    "",
    "二是季节性风险突出。冬季面临冰面安全和设施冻损风险；夏季汛期来临，山洪、滑坡、河道行洪等风险加剧",
    "",
    "三是工程类型复杂多样。既有水库、水电站等运行类工程，又有河道治理、水利设施提质增效等建设类项目",
    "",
    "四是与生态安全紧密关联。作为京津唐重要水源地，我市水务安全生产不仅关系人民群众生命财产安全，还直接影响区域生态环境和供水安全"
], font_size=15)

# 第 8 页：二 已开展的工作
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "已开展的工作")
add_content(slide, [
    "一是安全生产责任体系逐步完善",
    "• 严格落实"三管三必须"要求",
    "• 制定《承德市水务局有关领导干部安全生产职责清单》《局机关各科室和局属各单位安全生产职责清单》《承德市水务局安全生产工作指导意见》",
    "• 构建从业务科室到局属单位、从监管单位到生产经营单位的责任闭环",
    "",
    "二是专项整治行动深入开展",
    "• 聚焦工程建设、运行管理等重点领域",
    "• 持续开展治本攻坚三年行动、百日攻坚、冬春集中攻坚、常态化隐患排查、"双盲"抽查等专项行动",
    "• 事故发生率呈稳中有降趋势",
    "",
    "三是安全生产基础工作不断夯实",
    "• 常态化开展安全警示教育和业务培训",
    "• 推进水利安全生产风险管控"六项机制"建设",
    "• 完善应急预案体系，组织开展应急演练"
], font_size=14)

# 第 9 页：二（四）存在的主要问题
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "（四）存在的主要问题")
add_content(slide, [
    "监管层面存在"四大短板"：",
    "1. 责任传导"上热下冷"。市局部署的安全要求在部分县级部门、基层水管单位落实不到位",
    "2. 协同管控合力不足。市、县之间、科室之间联动机制不健全",
    "3. 执法监管刚性不够。部分检查流于形式，多以整改代替处罚",
    "4. 隐患识别能力偏弱。基层监管人员专业素养不足",
    "",
    "执行层面存在"三大问题"：",
    "1. 主体责任虚化。部分单位将安全管理视为附加任务，岗位责任未落实到人",
    "2. 隐患排查整治不彻底。对排查出的问题整改不及时、不彻底，存在"临时应付"心态",
    "3. 人员安全素养不足。部分从业人员安全意识淡薄，违规操作现象时有发生"
], font_size=14)

# 第 10 页：三 面临的形势
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "三、面临的形势")
add_content(slide, [
    "（一）工作标准越来越严",
    "国家对省级的安全生产考核改为由党中央负责组织实施，考核规格、标准大幅提升",
    "",
    "（二）风险考验越来越大",
    "近年来，极端天气事件呈现趋多趋频趋强趋广态势，加之降雨带"北抬"，安全生产与自然灾害风险交织叠加",
    "",
    "（三）事故发生概率增加",
    "当前灾后重建、国债支持项目、全市重点水利工程集中开工，高边坡、深基坑等高危作业增多",
    "",
    "（四）舆情关注度不断提高",
    "水务安全生产事关民生福祉，社会关注度高，一旦发生安全事故，易引发负面舆情"
], font_size=16)

# 第 11 页：四 下一步重点工作
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "四、下一步重点工作")
add_content(slide, [
    "一是持续健全责任体系",
    "动态调整职责清单，根据工作实际补充完善责任内容、明确责任边界，构建全链条、全流程、全覆盖的闭环管理体系",
    "",
    "二是持续深化专项整治",
    "完成《全市水利系统安全生产治本攻坚三年行动方案（2024-2026 年）》任务，对重点领域开展拉网式、全覆盖排查",
    "",
    "三是持续提升防控能力",
    "稳步推进水利安全生产风险管控"六项机制"建设，常态化开展安全警示教育和业务培训",
    "",
    "四是持续完善应急保障",
    "修订完善各类突发事件应急预案，组织常态化开展应急演练，通过实战化演练检验预案、锻炼队伍、提升能力"
], font_size=15)

# 第 12 页：五（一）工程施工安全隐患
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "五、常见安全隐患 - 工程施工环节")
add_content(slide, [
    "（一）水利工程施工环节六大安全隐患",
    "",
    "1. 现场管理不规范。无专项施工方案、不按方案施工、特种作业人员无证上岗",
    "",
    "2. 基坑、边坡、围堰风险突出。边坡放坡不符合要求、支护措施不到位、基坑排水不通畅",
    "",
    "3. 高处作业与脚手架、模板隐患。脚手架搭设不规范、模板支撑体系违规搭设、高处作业未系安全带",
    "",
    "4. 施工用电不规范。电线破损拖地泡水私拉乱接、漏电保护器失效",
    "",
    "5. 起重机械及施工设备隐患。未检测验收、限位装置失效、带病运行",
    "",
    "6. 消防安全管控不到位。易燃材料乱堆乱放、动火作业无审批、消防器材缺失失效"
], font_size=13)

# 第 13 页：五（二）设施运行安全隐患
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "五、常见安全隐患 - 设施运行环节")
add_content(slide, [
    "（二）水利设施运行环节五大安全隐患",
    "",
    "1. 水库、大坝安全隐患。坝体渗漏管涌裂缝滑坡、溢洪道堵塞、监测设施异常",
    "",
    "2. 水闸、泵站安全隐患。闸门变形锈蚀启闭失灵、机电设备老化、违规操作",
    "",
    "3. 堤防、河道安全隐患。堤防塌陷洞穴、河道淤积行洪断面不足、违章侵占河道",
    "",
    "4. 供水、灌区工程安全隐患。管道破损渗漏、阀门失灵、水源地污染隐患",
    "",
    "5. 管理与应急保障隐患。值班值守制度不落实、应急物资不足、应急演练流于形式"
], font_size=15)

# 第 14 页：结束页
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_PRIMARY
bg.line.fill.background()

text = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), prs.slide_width - Inches(1.6), Inches(2))
tf = text.text_frame
p = tf.paragraphs[0]
p.text = "汇报完毕"
p.alignment = PP_ALIGN.CENTER
set_font(p, size=36, bold=True, color=RGBColor(255, 255, 255))

p = tf.add_paragraph()
p.text = "请各位领导批评指正"
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(20)
set_font(p, size=24, color=RGBColor(255, 255, 255))

# 保存
output_path = "/Users/jiyi/.openclaw/workspace/水库安全生产汇报.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print("共 14 页，内容基于原文")
