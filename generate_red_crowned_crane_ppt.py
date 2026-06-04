#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
丹顶鹤冰箱贴设计方案 PPT - 深化设计版
基于设计方案文档进行信息可视化和逻辑梳理
参考水务安全生产汇报 PPT 的专业风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============ 配色方案（参考水务 PPT + 丹顶鹤主题） ============
COLOR_PRIMARY = RGBColor(25, 55, 109)      # 深蓝 - 主色
COLOR_SECONDARY = RGBColor(68, 114, 196)    # 中蓝 - 辅助色
COLOR_ACCENT = RGBColor(180, 30, 30)        # 丹顶红 - 强调色
COLOR_GOLD = RGBColor(218, 165, 32)         # 金色 - 点缀
COLOR_LIGHT = RGBColor(240, 245, 250)       # 浅蓝背景
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT = RGBColor(0, 0, 0)
COLOR_TEXT_LIGHT = RGBColor(80, 80, 80)
COLOR_CRANE_WHITE = RGBColor(248, 248, 248) # 鹤白色
COLOR_CRANE_BLACK = RGBColor(30, 30, 30)    # 鹤黑色

def set_font(paragraph, size=18, bold=False, color=COLOR_TEXT, font_name='微软雅黑'):
    """设置字体样式"""
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

def add_title_bar(slide, title, subtitle=None):
    """添加标题栏"""
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(10), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_font(p, size=24, bold=True, color=COLOR_WHITE)
    
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        set_font(p, size=14, color=RGBColor(200, 200, 200))

def add_image_with_border(slide, image_path, left, top, width, height, caption=None):
    """添加带边框和说明的图片"""
    try:
        pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        # 添加边框
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        border.fill.background()
        border.line.color.rgb = COLOR_SECONDARY
        border.line.width = Pt(2)
        
        if caption:
            cap_box = slide.shapes.add_textbox(left, top + height + Inches(0.1), width, Inches(0.4))
            tf = cap_box.text_frame
            p = tf.paragraphs[0]
            p.text = caption
            p.alignment = PP_ALIGN.CENTER
            set_font(p, size=11, color=COLOR_TEXT_LIGHT)
        return pic
    except Exception as e:
        # 如果图片不存在，创建占位符
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(245, 245, 245)
        placeholder.line.color.rgb = COLOR_SECONDARY
        placeholder.line.width = Pt(2)
        placeholder.line.dash_style = 4
        
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        p.text = f"【图片】\n{caption or image_path}"
        p.alignment = PP_ALIGN.CENTER
        set_font(p, size=12, color=COLOR_TEXT_LIGHT)
        return placeholder

def add_info_box(slide, title, content, left, top, width, height, bg_color=COLOR_LIGHT, title_color=COLOR_PRIMARY):
    """添加信息框"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = title_color
    box.line.width = Pt(1.5)
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    set_font(p, size=16, bold=True, color=title_color)
    
    p = tf.add_paragraph()
    p.text = "\n" + content
    p.space_after = Pt(8)
    set_font(p, size=13, color=COLOR_TEXT)
    return box

def add_comparison_table(slide, title, headers, rows, left, top, width, height):
    """添加对比表格"""
    # 表格背景
    table_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    table_bg.fill.solid()
    table_bg.fill.fore_color.rgb = COLOR_WHITE
    table_bg.line.color.rgb = COLOR_SECONDARY
    table_bg.line.width = Pt(2)
    
    # 创建表格
    rows_count = len(rows) + 1
    cols_count = len(headers)
    table = slide.shapes.add_table(rows_count, cols_count, left + Inches(0.2), top + Inches(0.3), 
                                   width - Inches(0.4), height - Inches(0.5)).table
    
    # 设置列宽
    col_width = int((width - Inches(0.4)) / cols_count)
    for i in range(cols_count):
        table.columns[i].width = col_width
    
    # 填充表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_font(p, size=14, bold=True, color=COLOR_WHITE)
    
    # 填充数据行
    colors = [RGBColor(248, 248, 248), RGBColor(255, 255, 255)]
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors[row_idx % 2]
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            set_font(p, size=12, color=COLOR_TEXT)

# ============ 开始制作 PPT ============

# 第 1 页：封面
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 背景渐变效果
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.5))
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_PRIMARY
bg.line.fill.background()

# 装饰性丹顶红圆点
red_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(0.3), Inches(1.2), Inches(1.2))
red_circle.fill.solid()
red_circle.fill.fore_color.rgb = COLOR_ACCENT
red_circle.line.fill.background()

# 标题
title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(10), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "丹顶鹤主题冰箱贴"
p.alignment = PP_ALIGN.LEFT
set_font(p, size=42, bold=True, color=COLOR_PRIMARY)

p = tf.add_paragraph()
p.text = "三款完整设计方案"
set_font(p, size=24, color=COLOR_TEXT_LIGHT)

# 副标题
sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(10), Inches(1))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "产品设计 · 视觉规划 · 工艺说明"
set_font(p, size=16, color=COLOR_SECONDARY)

# 版本信息
version_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(5), Inches(0.5))
tf = version_box.text_frame
p = tf.paragraphs[0]
p.text = "版本 V1.0 | 2026 年 3 月"
set_font(p, size=12, color=COLOR_TEXT_LIGHT)

# 第 2 页：产品概览
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "产品概览", "三款不同风格定位，覆盖多元消费人群")

# 三个产品卡片
products = [
    {
        "name": "松间仙鹤",
        "style": "国风浮雕款",
        "position": "高端文创",
        "target": "文化爱好者、游客",
        "price": "¥68-88",
        "color": RGBColor(25, 55, 109)
    },
    {
        "name": "萌鹤团子",
        "style": "可爱爆款款",
        "position": "大众爆款",
        "target": "年轻人、学生",
        "price": "¥35-45",
        "color": RGBColor(237, 125, 49)
    },
    {
        "name": "极线飞鹤",
        "style": "现代家居款",
        "position": "设计师款",
        "target": "品质生活人群",
        "price": "¥58-78",
        "color": RGBColor(30, 30, 30)
    }
]

for i, prod in enumerate(products):
    left = Inches(0.8) + i * Inches(4.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_LIGHT
    box.line.color.rgb = prod["color"]
    box.line.width = Pt(3)
    
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = prod["name"]
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=20, bold=True, color=prod["color"])
    
    p = tf.add_paragraph()
    p.text = f"\n{prod['style']}\n\n定位：{prod['position']}\n\n人群：{prod['target']}\n\n价格：{prod['price']}"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=14, color=COLOR_TEXT)

# 第 3 页：设计理念与逻辑
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "设计理念", "从传统文化到现代审美的完整表达")

# 设计理念图示
concept_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5))
concept_left.fill.solid()
concept_left.fill.fore_color.rgb = COLOR_LIGHT
concept_left.line.color.rgb = COLOR_PRIMARY
concept_left.line.width = Pt(2)

tf = concept_left.text_frame
p = tf.paragraphs[0]
p.text = "核心设计逻辑"
set_font(p, size=18, bold=True, color=COLOR_PRIMARY)

concepts = [
    "\n1. 文化传承",
    "   丹顶鹤在中国文化中象征长寿、吉祥、高雅",
    "   提取传统元素：松、鹤、日轮、祥云",
    "\n2. 分层表达",
    "   三款产品对应三种审美层次",
    "   传统 → 萌趣 → 极简",
    "\n3. 场景覆盖",
    "   景区文创 → 年轻礼品 → 家居装饰",
    "   全渠道、全人群覆盖"
]
for c in concepts:
    p = tf.add_paragraph()
    p.text = c
    set_font(p, size=13, color=COLOR_TEXT)

# 右侧关系图
# 绘制三层金字塔结构
pyramid_levels = [
    ("高端收藏", "松间仙鹤", "文化深度", RGBColor(25, 55, 109)),
    ("大众消费", "萌鹤团子", "情感共鸣", RGBColor(237, 125, 49)),
    ("品质生活", "极线飞鹤", "审美表达", RGBColor(30, 30, 30))
]

for i, (level, name, value, color) in enumerate(pyramid_levels):
    height = Inches(1.3)
    top = Inches(1.8) + i * (height + Inches(0.2))
    
    # 左侧标签
    label = slide.shapes.add_textbox(Inches(7), top, Inches(2), height)
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = level
    p.alignment = PP_ALIGN.RIGHT
    set_font(p, size=14, color=COLOR_TEXT_LIGHT)
    
    # 中间条形
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), top, Inches(3), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    
    tf = bar.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=14, bold=True, color=COLOR_WHITE)
    
    # 右侧价值
    value_label = slide.shapes.add_textbox(Inches(12.4), top, Inches(0.8), height)
    tf = value_label.text_frame
    p = tf.paragraphs[0]
    p.text = value
    set_font(p, size=12, color=COLOR_TEXT)

# 第 4 页：款式一详细设计 - 松间仙鹤
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "款式一：松间仙鹤", "国风浮雕款 · 东方文创旗舰")

# 左侧设计说明
design_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(6), Inches(3.5))
design_box.fill.solid()
design_box.fill.fore_color.rgb = COLOR_LIGHT
design_box.line.color.rgb = COLOR_PRIMARY

tf = design_box.text_frame
p = tf.paragraphs[0]
p.text = "视觉设计要点"
set_font(p, size=16, bold=True, color=COLOR_PRIMARY)

design_points = [
    "\n• 轮廓：竖版椭圆形，中式窗棂造型",
    "• 构图：丹顶鹤居中偏左，S 型鹤颈",
    "• 层次：4 层浮雕，高低错落",
    "• 元素：鹤、日轮、迎客松、远山、祥云"
]
for pt in design_points:
    p = tf.add_paragraph()
    p.text = pt
    set_font(p, size=13, color=COLOR_TEXT)

# 配色方案小色卡
color_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.2), Inches(6), Inches(1.8))
color_box.fill.solid()
color_box.fill.fore_color.rgb = COLOR_WHITE
color_box.line.color.rgb = COLOR_SECONDARY

tf = color_box.text_frame
p = tf.paragraphs[0]
p.text = "配色方案"
set_font(p, size=14, bold=True, color=COLOR_PRIMARY)

# 色卡展示
colors_info = [
    ("象牙白", RGBColor(245, 245, 240), "鹤身"),
    ("朱砂红", RGBColor(180, 30, 30), "鹤顶"),
    ("墨黑", RGBColor(30, 30, 30), "鹤翼"),
    ("青黛蓝绿", RGBColor(60, 100, 100), "山体"),
    ("香槟金", RGBColor(218, 165, 32), "日轮")
]

for i, (name, color, usage) in enumerate(colors_info):
    color_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2) + i * Inches(1.1), Inches(5.8), Inches(0.7), Inches(0.7))
    color_circle.fill.solid()
    color_circle.fill.fore_color.rgb = color
    color_circle.line.color.rgb = COLOR_TEXT
    color_circle.line.width = Pt(1)
    
    label = slide.shapes.add_textbox(Inches(1.2) + i * Inches(1.1), Inches(6.55), Inches(1), Inches(0.3))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=9, color=COLOR_TEXT)

# 右侧图片
add_image_with_border(slide, "/Users/jiyi/.openclaw/workspace/丹顶鹤冰箱贴设计图.png", 
                     Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5),
                     "丹顶鹤冰箱贴设计图")

# 第 5 页：款式一工艺与结构
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "松间仙鹤", "工艺说明与分层结构")

# 分层结构图示
layer_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(3), Inches(0.5))
tf = layer_title.text_frame
p = tf.paragraphs[0]
p.text = "分层结构（从上到下）"
set_font(p, size=16, bold=True, color=COLOR_PRIMARY)

layers = [
    ("第 1 层：鹤头层", "1.5mm 软 PVC", "最高浮雕，含鹤顶红点"),
    ("第 2 层：鹤身层", "1mm 软 PVC", "身体主体，羽毛纹理"),
    ("第 3 层：山石层", "1mm 软 PVC", "山石、松枝"),
    ("第 4 层：底板", "1.5mm 软 PVC", "背景、云雾、日轮"),
    ("第 5 层：磁片", "1mm 橡胶磁", "整面背胶")
]

for i, (name, spec, desc) in enumerate(layers):
    top = Inches(2) + i * Inches(0.9)
    
    # 层块
    layer_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(6), Inches(0.75))
    layer_box.fill.solid()
    layer_box.fill.fore_color.rgb = RGBColor(240, 240, 240) if i % 2 == 0 else RGBColor(250, 250, 250)
    layer_box.line.color.rgb = COLOR_SECONDARY
    
    tf = layer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{name}  |  {spec}  |  {desc}"
    set_font(p, size=12, color=COLOR_TEXT)

# 工艺说明
craft_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.3), Inches(5.3), Inches(3.5))
craft_box.fill.solid()
craft_box.fill.fore_color.rgb = COLOR_LIGHT
craft_box.line.color.rgb = COLOR_PRIMARY

tf = craft_box.text_frame
p = tf.paragraphs[0]
p.text = "核心工艺"
set_font(p, size=16, bold=True, color=COLOR_PRIMARY)

crafts = [
    "\n✓ 多层浮雕 - 4 层高低错落",
    "✓ 烫金工艺 - 日轮香槟金烫金",
    "✓ 半透明喷涂 - 云雾朦胧感",
    "✓ 细纹模具 - 羽毛纹理清晰",
    "✓ 渐变喷涂 - 墨黑到深灰过渡"
]
for c in crafts:
    p = tf.add_paragraph()
    p.text = c
    set_font(p, size=13, color=COLOR_TEXT)

# 尺寸规格
spec_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(5.2), Inches(5.3), Inches(1.8))
spec_box.fill.solid()
spec_box.fill.fore_color.rgb = COLOR_WHITE
spec_box.line.color.rgb = COLOR_ACCENT

tf = spec_box.text_frame
p = tf.paragraphs[0]
p.text = "尺寸规格：78×55×5mm | 重量：约 25g"
set_font(p, size=14, bold=True, color=COLOR_ACCENT)

# 第 6 页：款式二详细设计 - 萌鹤团子
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "款式二：萌鹤团子", "可爱爆款款 · 年轻礼品市场")

# 左侧设计特点
features_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(6), Inches(3.5))
features_box.fill.solid()
features_box.fill.fore_color.rgb = RGBColor(255, 245, 240)
features_box.line.color.rgb = RGBColor(237, 125, 49)

tf = features_box.text_frame
p = tf.paragraphs[0]
p.text = "萌系设计要素"
set_font(p, size=16, bold=True, color=RGBColor(180, 80, 0))

features = [
    "\n• 比例：大头小身 2:1，Q 版标准",
    "• 表情：闭眼笑 + 弯弯眼，治愈感",
    "• 腮红：两团蜜桃粉，增加可爱度",
    "• 动作：双手抱小鱼，增加故事性",
    "• 背景：水波纹 + 气泡 + 荷叶，场景感"
]
for f in features:
    p = tf.add_paragraph()
    p.text = f
    set_font(p, size=13, color=COLOR_TEXT)

# 配色方案
color_palette = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.2), Inches(6), Inches(1.8))
color_palette.fill.solid()
color_palette.fill.fore_color.rgb = COLOR_WHITE
color_palette.line.color.rgb = RGBColor(237, 125, 49)

tf = color_palette.text_frame
p = tf.paragraphs[0]
p.text = "马卡龙配色方案"
set_font(p, size=14, bold=True, color=RGBColor(180, 80, 0))

cute_colors = [
    ("纯白", RGBColor(255, 255, 255), "鹤身"),
    ("鲜红", RGBColor(220, 60, 60), "鹤顶"),
    ("蜜桃粉", RGBColor(255, 180, 180), "腮红"),
    ("马卡龙蓝", RGBColor(180, 220, 240), "水波"),
    ("嫩绿", RGBColor(180, 230, 180), "荷叶")
]

for i, (name, color, usage) in enumerate(cute_colors):
    color_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2) + i * Inches(1.1), Inches(5.8), Inches(0.7), Inches(0.7))
    color_circle.fill.solid()
    color_circle.fill.fore_color.rgb = color
    color_circle.line.color.rgb = COLOR_TEXT
    color_circle.line.width = Pt(1)
    
    label = slide.shapes.add_textbox(Inches(1.2) + i * Inches(1.1), Inches(6.55), Inches(1), Inches(0.3))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=9, color=COLOR_TEXT)

# 右侧图片占位
add_image_with_border(slide, "/Users/jiyi/.openclaw/workspace/丹顶鹤冰箱贴设计图.png",
                     Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5),
                     "萌鹤团子设计效果图")

# 第 7 页：款式二工艺与材质
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "萌鹤团子", "滴胶亚克力工艺")

# 材质结构
structure_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(6), Inches(3.5))
structure_box.fill.solid()
structure_box.fill.fore_color.rgb = COLOR_LIGHT
structure_box.line.color.rgb = COLOR_SECONDARY

tf = structure_box.text_frame
p = tf.paragraphs[0]
p.text = "分层结构（从上到下）"
set_font(p, size=16, bold=True, color=COLOR_PRIMARY)

layers2 = [
    ("滴胶层", "1.5mm UV 树脂", "高透明，表面光滑，水晶质感"),
    ("印刷层", "0.1mm 彩印", "1200dpi 高清印刷"),
    ("底板", "2mm 亚克力", "透明或白色底板"),
    ("磁片", "0.8mm 橡胶磁", "圆形背胶")
]

for i, (name, spec, desc) in enumerate(layers2):
    top = Inches(2) + i * Inches(0.8)
    layer_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(6), Inches(0.65))
    layer_box.fill.solid()
    layer_box.fill.fore_color.rgb = RGBColor(245, 245, 245) if i % 2 == 0 else RGBColor(255, 255, 255)
    layer_box.line.color.rgb = COLOR_SECONDARY
    
    tf = layer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{name}  |  {spec}  |  {desc}"
    set_font(p, size=12, color=COLOR_TEXT)

# 特色工艺
craft_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.3), Inches(5.3), Inches(2.8))
craft_box.fill.solid()
craft_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
craft_box.line.color.rgb = RGBColor(237, 125, 49)

tf = craft_box.text_frame
p = tf.paragraphs[0]
p.text = "特色工艺"
set_font(p, size=16, bold=True, color=RGBColor(180, 80, 0))

crafts2 = [
    "\n✓ 滴胶 - 高透明，镜面效果",
    "✓ 高清彩印 - 1200dpi，色彩鲜艳",
    "✓ 可选：微闪粉、夜光效果"
]
for c in crafts2:
    p = tf.add_paragraph()
    p.text = c
    set_font(p, size=13, color=COLOR_TEXT)

# 尺寸规格
spec_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(4.5), Inches(5.3), Inches(1.5))
spec_box.fill.solid()
spec_box.fill.fore_color.rgb = COLOR_WHITE
spec_box.line.color.rgb = RGBColor(237, 125, 49)

tf = spec_box.text_frame
p = tf.paragraphs[0]
p.text = "尺寸规格：Φ60×4mm | 重量：约 18g"
set_font(p, size=14, bold=True, color=RGBColor(180, 80, 0))

# 第 8 页：款式三详细设计 - 极线飞鹤
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "款式三：极线飞鹤", "现代家居款 · 北欧风极简主义")

# 设计理念
concept_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(6), Inches(2.5))
concept_box.fill.solid()
concept_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
concept_box.line.color.rgb = COLOR_CRANE_BLACK

tf = concept_box.text_frame
p = tf.paragraphs[0]
p.text = "极简设计理念"
set_font(p, size=16, bold=True, color=COLOR_CRANE_BLACK)

concept_text = [
    "\n• 造型：线条化，几何切面表达翅膀",
    "• 构图：一条连续弧线完成身体",
    "• 留白：占画面 70%，营造空间感",
    "• 元素：仅保留鹤、水面线、红圆"
]
for t in concept_text:
    p = tf.add_paragraph()
    p.text = t
    set_font(p, size=13, color=COLOR_TEXT)

# 配色方案（极简）
minimal_colors = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.2), Inches(6), Inches(2.8))
minimal_colors.fill.solid()
minimal_colors.fill.fore_color.rgb = COLOR_WHITE
minimal_colors.line.color.rgb = COLOR_CRANE_BLACK

tf = minimal_colors.text_frame
p = tf.paragraphs[0]
p.text = "极简双色 + 点缀"
set_font(p, size=14, bold=True, color=COLOR_CRANE_BLACK)

p = tf.add_paragraph()
p.text = "\n配色比例：\n底色（留白）85% + 线条 12% + 红圆 3%"
set_font(p, size=12, color=COLOR_TEXT)

# 三个色块
color_blocks = [
    ("哑光奶油白", RGBColor(250, 248, 240), "85%"),
    ("深炭黑", RGBColor(30, 30, 30), "12%"),
    ("正红", RGBColor(180, 30, 30), "3%")
]

for i, (name, color, ratio) in enumerate(color_blocks):
    color_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2) + i * Inches(2), Inches(5.5), Inches(1.5), Inches(1))
    color_rect.fill.solid()
    color_rect.fill.fore_color.rgb = color
    color_rect.line.color.rgb = COLOR_TEXT
    color_rect.line.width = Pt(1)
    
    label = slide.shapes.add_textbox(Inches(1.2) + i * Inches(2), Inches(6.6), Inches(1.5), Inches(0.3))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = f"{name}\n{ratio}"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=10, color=COLOR_TEXT)

# 右侧图片
add_image_with_border(slide, "/Users/jiyi/.openclaw/workspace/丹顶鹤冰箱贴设计图.png",
                     Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5),
                     "极线飞鹤设计效果图")

# 第 9 页：款式三工艺与材质
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "极线飞鹤", "陶瓷 + 哑光釉工艺")

# 工艺流程
process_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(6), Inches(3.5))
process_box.fill.solid()
process_box.fill.fore_color.rgb = COLOR_LIGHT
process_box.line.color.rgb = COLOR_CRANE_BLACK

tf = process_box.text_frame
p = tf.paragraphs[0]
p.text = "陶瓷工艺流程"
set_font(p, size=16, bold=True, color=COLOR_CRANE_BLACK)

process_steps = [
    "\n1. 瓷胚高温烧制成型",
    "2. 激光雕刻鹤线条（精细、流畅）",
    "3. 手工填充釉料（每只略有不同）",
    "4. 二次低温烧制",
    "5. 背面粘贴强磁片（钕磁铁）"
]
for step in process_steps:
    p = tf.add_paragraph()
    p.text = step
    set_font(p, size=13, color=COLOR_TEXT)

# 工艺特点
craft_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.3), Inches(5.3), Inches(2.8))
craft_box.fill.solid()
craft_box.fill.fore_color.rgb = RGBColor(245, 245, 240)
craft_box.line.color.rgb = COLOR_CRANE_BLACK

tf = craft_box.text_frame
p = tf.paragraphs[0]
p.text = "核心工艺"
set_font(p, size=16, bold=True, color=COLOR_CRANE_BLACK)

crafts3 = [
    "\n✓ 激光刻线 - 精细流畅，一致性好",
    "✓ 手工上釉 - 独特性，每只不同",
    "✓ 哑光釉 - 温润质感，不反光",
    "✓ 强磁片 - 钕磁铁，吸附力强"
]
for c in crafts3:
    p = tf.add_paragraph()
    p.text = c
    set_font(p, size=13, color=COLOR_TEXT)

# 尺寸规格
spec_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(4.5), Inches(5.3), Inches(1.5))
spec_box.fill.solid()
spec_box.fill.fore_color.rgb = COLOR_WHITE
spec_box.line.color.rgb = COLOR_CRANE_BLACK

tf = spec_box.text_frame
p = tf.paragraphs[0]
p.text = "尺寸规格：85×40×3.5mm | 重量：约 35g"
set_font(p, size=14, bold=True, color=COLOR_CRANE_BLACK)

# 第 10 页：三款产品对比分析
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "三款产品对比", "定位、材质、工艺、渠道全方位对比")

# 对比表格
headers = ["对比项", "松间仙鹤", "萌鹤团子", "极线飞鹤"]
rows = [
    ["风格定位", "国风文创", "Q 版可爱", "现代极简"],
    ["目标人群", "文化爱好者/游客", "年轻人/学生", "品质生活人群"],
    ["材质", "软 PVC 浮雕", "滴胶亚克力", "陶瓷"],
    ["尺寸", "78×55×5mm", "Φ60×4mm", "85×40×3.5mm"],
    ["重量", "~25g", "~18g", "~35g"],
    ["工艺复杂度", "★★★☆☆", "★★☆☆☆", "★★★★☆"],
    ["建议零售价", "¥68-88", "¥35-45", "¥58-78"],
    ["目标渠道", "景区/博物馆", "电商/市集", "买手店/联名"],
    ["生产周期", "15-20 天", "10-15 天", "25-30 天"],
    ["起订量", "500 个", "300 个", "200 个"]
]

add_comparison_table(slide, "", headers, rows, Inches(0.8), Inches(1.3), Inches(11.8), Inches(5.5))

# 第 11 页：成本与定价策略
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "成本估算与定价", "基于 1000 个起订量估算")

# 成本对比柱状图（用矩形模拟）
cost_data = [
    ("松间仙鹤", 17, 24, 68, 88, RGBColor(25, 55, 109)),
    ("萌鹤团子", 11, 17, 35, 45, RGBColor(237, 125, 49)),
    ("极线飞鹤", 32, 46, 58, 78, RGBColor(30, 30, 30))
]

for i, (name, cost_min, cost_max, price_min, price_max, color) in enumerate(cost_data):
    top = Inches(1.5) + i * Inches(2)
    
    # 产品名称
    name_box = slide.shapes.add_textbox(Inches(0.8), top, Inches(2), Inches(0.5))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = name
    set_font(p, size=14, bold=True, color=color)
    
    # 成本条
    cost_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), top, Inches(2.5), Inches(0.5))
    cost_bar.fill.solid()
    cost_bar.fill.fore_color.rgb = RGBColor(200, 200, 200)
    cost_bar.line.fill.background()
    
    tf = cost_bar.text_frame
    p = tf.paragraphs[0]
    p.text = f"成本 ¥{cost_min}-{cost_max}"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=11, color=COLOR_WHITE)
    
    # 售价条
    price_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.7), top, Inches(4), Inches(0.5))
    price_bar.fill.solid()
    price_bar.fill.fore_color.rgb = color
    price_bar.line.fill.background()
    
    tf = price_bar.text_frame
    p = tf.paragraphs[0]
    p.text = f"售价 ¥{price_min}-{price_max}"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=11, color=COLOR_WHITE)
    
    # 毛利率
    margin = int((1 - (cost_min + cost_max) / 2 / ((price_min + price_max) / 2)) * 100)
    margin_box = slide.shapes.add_textbox(Inches(10), top, Inches(2.5), Inches(0.5))
    tf = margin_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"毛利率 ~{margin}%"
    set_font(p, size=14, bold=True, color=COLOR_ACCENT)

# 成本明细
cost_detail = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2), Inches(11.8), Inches(1.8))
cost_detail.fill.solid()
cost_detail.fill.fore_color.rgb = COLOR_LIGHT
cost_detail.line.color.rgb = COLOR_SECONDARY

tf = cost_detail.text_frame
p = tf.paragraphs[0]
p.text = "成本构成（以松间仙鹤为例）"
set_font(p, size=14, bold=True, color=COLOR_PRIMARY)

p = tf.add_paragraph()
p.text = "\n材料成本 ¥8-12 (47%)  |  加工成本 ¥6-8 (35%)  |  包装成本 ¥3-4 (18%)"
set_font(p, size=13, color=COLOR_TEXT)

# 第 12 页：打样与生产流程
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "打样与生产", "从设计到交付的完整流程")

# 流程图
flow_steps = [
    ("设计确认", "矢量文件 AI/CDR\n3D 文件 STP/IGS\n效果图 PNG/JPG", Inches(0.8), RGBColor(25, 55, 109)),
    ("打样制作", "3-5 个样品\n7-15 工作日\n费用可退还", Inches(4), RGBColor(237, 125, 49)),
    ("样品确认", "颜色对色卡\n尺寸测量\n工艺确认", Inches(7.2), RGBColor(68, 114, 196)),
    ("批量生产", "按订单量生产\n15-30 天周期\n质检包装", Inches(10.4), RGBColor(30, 30, 30))
]

for i, (title, desc, left, color) in enumerate(flow_steps):
    # 步骤框
    step_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2), Inches(2.8), Inches(2.5))
    step_box.fill.solid()
    step_box.fill.fore_color.rgb = COLOR_LIGHT
    step_box.line.color.rgb = color
    step_box.line.width = Pt(2.5)
    
    tf = step_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=15, bold=True, color=color)
    
    p = tf.add_paragraph()
    p.text = "\n" + desc
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=11, color=COLOR_TEXT)
    
    # 箭头（除了最后一个）
    if i < len(flow_steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(2.9), Inches(2.8), Inches(0.8), Inches(0.8))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_SECONDARY
        arrow.line.fill.background()

# 打样要求
req_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5), Inches(11.8), Inches(2))
req_box.fill.solid()
req_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
req_box.line.color.rgb = COLOR_ACCENT

tf = req_box.text_frame
p = tf.paragraphs[0]
p.text = "打样通用要求"
set_font(p, size=14, bold=True, color=COLOR_ACCENT)

p = tf.add_paragraph()
p.text = "\n• 打样数量：每款 3-5 个  |  • 打样周期：7-15 工作日  |  • 打样费用：大货下单后可退还"
p = tf.add_paragraph()
p.text = "• 颜色确认：提供 Pantone 色卡对色  |  • 尺寸确认：游标卡尺测量，公差范围内"
set_font(p, size=12, color=COLOR_TEXT)

# 第 13 页：质量标准与认证
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "质量标准", "符合国际环保与安全标准")

# 三个产品的质量标准
standards = [
    {
        "name": "松间仙鹤 (PVC)",
        "standards": [
            "色牢度 ≥4 级 (GB/T 251)",
            "附着力 ≥4B (GB/T 9286)",
            "耐温性 -20℃~60℃",
            "环保 EN71-3、REACH"
        ],
        "color": RGBColor(25, 55, 109)
    },
    {
        "name": "萌鹤团子 (滴胶)",
        "standards": [
            "滴胶硬度 ≥80 Shore D",
            "透光率 ≥90%",
            "耐黄变 ≥1000 小时 (UV)",
            "环保 EN71-3、REACH"
        ],
        "color": RGBColor(237, 125, 49)
    },
    {
        "name": "极线飞鹤 (陶瓷)",
        "standards": [
            "吸水率 ≤0.5%",
            "抗折强度 ≥35MPa",
            "釉面硬度 ≥6 Mohs",
            "耐温 -10℃~100℃, FDA、LFGB"
        ],
        "color": RGBColor(30, 30, 30)
    }
]

for i, std in enumerate(standards):
    left = Inches(0.8) + i * Inches(4.2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.8), Inches(5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_LIGHT
    box.line.color.rgb = std["color"]
    box.line.width = Pt(2.5)
    
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = std["name"]
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=14, bold=True, color=std["color"])
    
    for s in std["standards"]:
        p = tf.add_paragraph()
        p.text = "\n" + s
        p.alignment = PP_ALIGN.LEFT
        set_font(p, size=11, color=COLOR_TEXT)

# 第 14 页：包装与物流
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "包装方案", "保护产品 + 提升体验")

# 三款产品包装对比
packaging = [
    {
        "name": "松间仙鹤",
        "inner": "透明 OPP 袋 + 卡纸背板",
        "card": "100×150mm, 350g 铜版纸覆哑膜",
        "outer": "瓦楞纸盒 10 个/盒",
        "carton": "五层纸箱 100 个/箱"
    },
    {
        "name": "萌鹤团子",
        "inner": "透明 PET 圆盒 + 彩卡",
        "card": "80×100mm, 300g 铜版纸覆亮膜",
        "outer": "瓦楞纸盒 20 个/盒",
        "carton": "五层纸箱 200 个/箱"
    },
    {
        "name": "极线飞鹤",
        "inner": "天地盖纸盒 + EVA 内托",
        "card": "100×60×20mm, 1200g 灰板特种纸",
        "outer": "瓦楞纸盒 10 个/盒",
        "carton": "五层纸箱 100 个/箱"
    }
]

for i, pack in enumerate(packaging):
    top = Inches(1.5) + i * Inches(1.8)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(11.8), Inches(1.6))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_LIGHT if i % 2 == 0 else COLOR_WHITE
    box.line.color.rgb = COLOR_SECONDARY
    
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{pack['name']}  |  内包装：{pack['inner']}  |  卡纸：{pack['card']}  |  外箱：{pack['carton']}"
    set_font(p, size=12, color=COLOR_TEXT)

# 第 15 页：封底
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_PRIMARY
bg.line.fill.background()

# 装饰性红圆
red_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(1), Inches(2), Inches(2))
red_circle.fill.solid()
red_circle.fill.fore_color.rgb = COLOR_ACCENT
red_circle.line.fill.background()

# 结束语
text_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), prs.slide_width - Inches(1.6), Inches(2))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = "感谢审阅"
p.alignment = PP_ALIGN.CENTER
set_font(p, size=36, bold=True, color=COLOR_WHITE)

p = tf.add_paragraph()
p.text = "\n欢迎洽谈合作"
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(20)
set_font(p, size=24, color=RGBColor(200, 200, 200))

# 联系信息
contact_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(6), Inches(1))
tf = contact_box.text_frame
p = tf.paragraphs[0]
p.text = "设计方：[您的公司/工作室名称]  |  版本 V1.0  |  2026 年 3 月"
p.alignment = PP_ALIGN.CENTER
set_font(p, size=12, color=RGBColor(180, 180, 180))

# 保存 PPT
output_path = "/Users/jiyi/.openclaw/workspace/丹顶鹤冰箱贴设计方案_完整版.pptx"
prs.save(output_path)
print(f"✅ PPT 已生成：{output_path}")
print(f"📊 共 15 页，包含：")
print("   - 产品概览与定位")
print("   - 三款产品详细设计（设计 + 工艺 + 材质）")
print("   - 对比分析表")
print("   - 成本估算与定价策略")
print("   - 生产流程与质量标准")
print("   - 包装方案")
print("\n💡 深化设计亮点：")
print("   - 信息可视化（对比表、流程图、色卡）")
print("   - 逻辑层次清晰（定位→设计→工艺→成本→生产）")
print("   - 专业配色与布局")
print("   - 图片与文字有机结合")
