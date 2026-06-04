#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""水务安全生产汇报 PPT 生成器 - 精细版"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLORS = {
    'primary': RGBColor(18, 90, 144),
    'secondary': RGBColor(64, 158, 204),
    'accent': RGBColor(237, 125, 49),
    'light': RGBColor(242, 247, 251),
    'dark': RGBColor(62, 62, 62),
    'white': RGBColor(255, 255, 255),
    'red': RGBColor(192, 80, 77),
}

def set_bg(slide, color=COLORS['light']):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    header.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(2.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    tf = tb2.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(26)
    p.font.color.rgb = COLORS['dark']
    p.alignment = PP_ALIGN.CENTER
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), prs.slide_height-Inches(0.3), Inches(5.333), Inches(0.3))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    return slide

def add_section_slide(num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    
    sb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.5), prs.slide_height)
    sb.fill.solid()
    sb.fill.fore_color.rgb = COLORS['primary']
    
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(2), Inches(4), Inches(3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(140)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    tb2 = slide.shapes.add_textbox(Inches(5.5), Inches(3), Inches(7), Inches(2))
    tf = tb2.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    return slide

def add_content_slide(title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(10), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    cb = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12), Inches(5.3))
    tf = cb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark']
        p.space_after = Pt(8)
    return slide

def add_stat_slide(title, stats):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['primary']
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(10), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    cw, ch = Inches(3.5), Inches(2.8)
    for i, stat in enumerate(stats):
        col, row = i % 3, i // 3
        x, y = Inches(1.2)+col*(cw+Inches(0.4)), Inches(2.3)+row*(ch+Inches(0.4))
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['primary']
        
        tb = slide.shapes.add_textbox(x+Inches(0.3), y+Inches(0.4), cw-Inches(0.6), Inches(1.2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = stat['n']
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        tb2 = slide.shapes.add_textbox(x+Inches(0.3), y+Inches(1.7), cw-Inches(0.6), Inches(0.9))
        tf = tb2.text_frame
        p = tf.paragraphs[0]
        p.text = stat['l']
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    return slide

def add_end_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    tf = tb2.text_frame
    p = tf.paragraphs[0]
    p.text = '谢谢观看'
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    return slide

# 生成 PPT
print('正在生成 PPT...')

add_title_slide('守牢安全底线 护航水务发展', '水务安全生产工作汇报')

add_content_slide('目录', [
    '01 切实加强水务安全生产工作的认识',
    '02 水务安全生产现状',
    '03 面临的形势',
    '04 下一步重点工作',
    '05 常见安全隐患'
])

add_section_slide('01', '切实加强水务安全生产工作的认识')

add_content_slide('书记关于安全生产重要论述核心要义', [
    '一是树牢红线意识，坚守发展决不能以牺牲人的生命为代价这条不可逾越的底线',
    '二是健全责任体系，严格落实党政同责、一岗双责、齐抓共管、失职追责',
    '三是压实主体责任，督促生产经营单位履行安全生产法定职责',
    '四是创新监管方式，加快推进安全监管改革，提升监管精准度和效能',
    '五是构建长效机制，实现风险常态化管控、隐患动态化治理',
    '六是强化责任担当，以时时放心不下的责任感抓细抓实安全生产各项工作'
])

add_content_slide('核心法律法规与政策文件', [
    '通用法规：《安全生产法》',
    '- 明确了企业主体责任、政府监管职责',
    '- 隐患排查治理和应急处置等基本原则',
    '',
    '水利专项法规：《水法》《防洪法》',
    '- 水利工程设计、施工、运行、维护全生命周期安全管理',
    '- 强调事前防控、过程控制、事后处置的闭环管理',
    '',
    '省市制度：《安全生产末端落实工作机制》等',
    '- 要求更严、标准更高、落点更实'
])

add_content_slide('水务安全生产目标', [
    '总体目标：',
    '- 坚决杜绝重特大事故',
    '- 遏制较大事故',
    '- 减少一般事故',
    '- 守牢安全底线',
    '',
    '基本底线：人员不伤亡、财产少损失',
    '',
    '确保：',
    '- 全市水利工程建设有序推进',
    '- 各类水利设施安全稳定运行'
])

add_section_slide('02', '水务安全生产现状')

add_stat_slide('2021-2023 年水利行业事故统计', [
    {'n': '136', 'l': '事故总数 (起)'},
    {'n': '152', 'l': '死亡人数 (人)'},
    {'n': '28', 'l': '受伤人数 (人)'},
    {'n': '68%', 'l': '建设项目事故'},
    {'n': '22%', 'l': '水库电站事故'}
])

add_content_slide('事故类型分布', [
    '水利建设项目事故 (68%)：',
    '- 高处坠落',
    '- 物体打击',
    '- 溺水',
    '- 机械伤害',
    '',
    '水库及水电站运行事故 (22%)：',
    '- 设施故障引发的溃坝',
    '- 触电等事故'
])

add_content_slide('水务安全生产四大特点', [
    '一是风险点多面广',
    '- 地处四河之源、两库上游',
    '- 四大水系、1500 条河流、众多水库和水电站',
    '- 安全监管覆盖范围广、难度大',
    '',
    '二是季节性风险突出',
    '- 冬季：冰面安全、设施冻损',
    '- 夏季：山洪、滑坡、河道行洪',
    '',
    '三是工程类型复杂多样',
    '- 运行类 + 建设类，老旧改造 + 新建施工',
    '',
    '四是与生态安全紧密关联',
    '- 京津唐重要水源地'
])

add_content_slide('已开展的工作', [
    '一是安全生产责任体系逐步完善',
    '- 严格落实三管三必须要求',
    '- 制定职责清单和指导意见',
    '- 构建责任闭环',
    '',
    '二是专项整治行动深入开展',
    '- 治本攻坚三年行动、百日攻坚等',
    '- 常态化隐患排查、双盲抽查',
    '',
    '三是安全生产基础工作不断夯实',
    '- 安全警示教育和业务培训',
    '- 推进六项机制建设',
    '- 完善应急预案体系'
])

add_content_slide('存在的主要问题', [
    '监管层面四大短板：',
    '- 责任传导上热下冷',
    '- 协同管控合力不足',
    '- 执法监管刚性不够',
    '- 隐患识别能力偏弱',
    '',
    '执行层面三大问题：',
    '- 主体责任虚化',
    '- 隐患排查整治不彻底',
    '- 人员安全素养不足'
])

add_section_slide('03', '面临的形势')

add_content_slide('当前面临的严峻形势', [
    '工作标准越来越严',
    '- 国家对省级考核由党中央组织实施',
    '- 较大事故国务院安委办挂牌督办',
    '',
    '风险考验越来越大',
    '- 极端天气趋多趋频趋强趋广',
    '- 安全与自然灾害风险交织叠加',
    '',
    '事故发生概率增加',
    '- 灾后重建、国债项目集中开工',
    '- 高危作业增多、人员流动性大',
    '',
    '舆情关注度不断提高',
    '- 社会关注度高，易引发负面舆情'
])

add_section_slide('04', '下一步重点工作')

add_content_slide('下一步重点工作', [
    '一是持续健全责任体系',
    '- 动态调整职责清单',
    '- 构建全链条、全流程、全覆盖闭环管理体系',
    '',
    '二是持续深化专项整治',
    '- 完成三年行动方案',
    '- 拉网式、全覆盖排查',
    '- 打非治违',
    '',
    '三是持续提升防控能力',
    '- 推进六项机制建设',
    '- 开展安全警示教育',
    '',
    '四是持续完善应急保障',
    '- 修订完善应急预案',
    '- 常态化开展应急演练'
])

add_section_slide('05', '常见安全隐患')

add_content_slide('水利工程施工环节六大安全隐患', [
    '一、现场管理不规范 - 无专项方案、特种作业无证上岗',
    '二、基坑边坡围堰风险 - 支护措施不到位、监测缺失',
    '三、高处作业与脚手架隐患 - 搭设不规范、安全带未系',
    '四、施工用电不规范 - 三级配电未执行、私拉乱接',
    '五、起重机械及设备隐患 - 未检测验收、带病运行',
    '六、消防安全管控不到位 - 动火无审批、器材缺失'
])

add_content_slide('水利设施运行环节五大安全隐患', [
    '一、水库大坝安全隐患',
    '- 坝体渗漏、管涌、裂缝',
    '- 溢洪道堵塞、监测设施异常',
    '',
    '二、水闸泵站安全隐患',
    '- 闸门变形、启闭失灵',
    '- 机电设备老化',
    '',
    '三、堤防河道安全隐患',
    '- 堤防塌陷、洞穴',
    '- 河道淤积、违章侵占',
    '',
    '四、供水灌区工程隐患',
    '- 管道破损、阀门失灵',
    '',
    '五、管理与应急保障隐患',
    '- 值班值守未落实、应急物资不足'
])

add_end_slide('安全生产重于泰山 责任落实贵在坚守')

output = '/Users/jiyi/.openclaw/workspace/水务安全生产汇报 - 精细版.pptx'
prs.save(output)
print(f'完成！文件：{output}')
