#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
水库安全生产汇报 PPT - 专业版
带图片占位符和图示区域
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLOR_PRIMARY = RGBColor(25, 55, 109)      # 深蓝
COLOR_SECONDARY = RGBColor(68, 114, 196)    # 中蓝
COLOR_ACCENT = RGBColor(237, 125, 49)       # 橙色
COLOR_LIGHT = RGBColor(240, 245, 250)       # 浅蓝背景
COLOR_TEXT = RGBColor(0, 0, 0)
COLOR_TEXT_LIGHT = RGBColor(80, 80, 80)

def set_font(paragraph, size=18, bold=False, color=COLOR_TEXT):
    for run in paragraph.runs:
        run.font.name = '微软雅黑'
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

def add_title_bar(slide, title):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), prs.slide_width - Inches(1.2), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_font(p, size=24, bold=True, color=RGBColor(255, 255, 255))

def add_image_placeholder(slide, left, top, width, height, label="插入图片"):
    """添加图片占位符"""
    placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(245, 245, 245)
    placeholder.line.color.rgb = COLOR_SECONDARY
    placeholder.line.width = Pt(2)
    placeholder.line.dash_style = 4  # 虚线
    
    tf = placeholder.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"【{label}】\n建议尺寸：{int(width*100)}x{int(height*100)}px\n点击插入相关图片"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=12, color=COLOR_TEXT_LIGHT)
    return placeholder

def add_content_with_image(slide, title, text_items, img_left, img_top, img_width, img_height, img_label):
    """添加带图片的内容页"""
    add_title_bar(slide, title)
    
    # 文本区域
    text_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(7), Inches(5.5))
    tf = text_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(text_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.space_after = Pt(8)
        set_font(p, size=15, color=COLOR_TEXT)
    
    # 图片占位符
    add_image_placeholder(slide, img_left, img_top, img_width, img_height, img_label)

def add_full_image_slide(slide, title, img_label):
    """全屏图片页"""
    add_title_bar(slide, title)
    add_image_placeholder(slide, Inches(0.6), Inches(1.2), prs.slide_width - Inches(1.2), Inches(5.5), img_label)

# ============ 开始制作 PPT ============

# 第 1 页：封面
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2))
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_PRIMARY
bg.line.fill.background()

# 背景图占位
add_image_placeholder(slide, Inches(8), Inches(0.5), Inches(4.5), Inches(3), "封面背景图\n水库/大坝全景照片")

title = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(6), Inches(1.5))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "守牢安全底线 护航水务发展"
p.alignment = PP_ALIGN.LEFT
set_font(p, size=38, bold=True, color=COLOR_PRIMARY)

sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(6), Inches(0.8))
tf = sub.text_frame
p = tf.paragraphs[0]
p.text = "承德市水务局安全生产工作汇报"
set_font(p, size=20, color=COLOR_TEXT_LIGHT)

# 第 2 页：目录
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "汇报目录")

# 左侧目录
dirs = [
    "01  切实加强水务安全生产\n      工作的认识",
    "02  水务安全生产现状",
    "03  面临的形势",
    "04  下一步重点工作",
    "05  常见安全隐患"
]
dir_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5), Inches(5))
tf = dir_box.text_frame
for i, d in enumerate(dirs):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = d
    p.space_after = Pt(20)
    set_font(p, size=20, color=COLOR_PRIMARY if i % 2 == 0 else COLOR_TEXT)

# 右侧图片
add_image_placeholder(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(5), "目录配图\n水务工作场景")

# 第 3 页：书记重要论述
content3 = [
    "（一）书记关于安全生产重要论述核心要义",
    "",
    "书记始终将安全生产作为极重的民生大事，纳入总体布局和战略布局统筹推进",
    "",
    "核心要牢牢把握六个方面：",
    "1. 树牢红线意识 - 发展决不能以牺牲人的生命为代价",
    "2. 健全责任体系 - 党政同责、一岗双责、齐抓共管、失职追责",
    "3. 压实主体责任 - 安全投入、管理、培训、防护、应急五到位",
    "4. 创新监管方式 - 提升监管精准度和效能",
    "5. 构建长效机制 - 风险常态化管控、隐患动态化治理",
    "6. 强化责任担当 - 以时时放心不下的责任感抓细抓实各项工作"
]
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_with_image(slide, "一、切实加强水务安全生产工作的认识", content3, Inches(8), Inches(1.3), Inches(4.5), Inches(5), "领导视察/会议照片")

# 第 4 页：法律法规
content4 = [
    "（二）核心法律法规与政策文件",
    "",
    "通用层面：",
    "• 《安全生产法》- 企业主体责任、政府监管职责、隐患排查治理",
    "",
    "水利领域：",
    "• 《水法》《防洪法》- 全生命周期安全管理",
    "• 事前防控、过程控制、事后处置闭环管理",
    "",
    "省市制度：",
    "• 《安全生产末端落实工作机制》",
    "• 《安全生产违法行为行政处罚规定》"
]
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_with_image(slide, "法律法规体系", content4, Inches(7.5), Inches(1.3), Inches(5), Inches(3), "法律法规封面图")

# 第 5 页：工作目标
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "（三）水务安全生产的目标")

# 目标框
targets = [
    ("省水利厅总体目标", "坚决杜绝重特大事故\n遏制较大事故\n减少一般事故\n守牢安全底线"),
    ("我市基本底线", "人员不伤亡\n财产少损失"),
    ("工作保障", "水利工程建设有序推进\n水利设施安全稳定运行\n为经济社会发展提供水安全保障")
]

for i, (title, text) in enumerate(targets):
    top = Inches(1.3) + i * Inches(2)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top, Inches(5.5), Inches(1.7))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_LIGHT if i % 2 == 0 else RGBColor(255, 250, 240)
    box.line.color.rgb = COLOR_SECONDARY
    
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_font(p, size=18, bold=True, color=COLOR_PRIMARY)
    p = tf.add_paragraph()
    p.text = text
    set_font(p, size=15, color=COLOR_TEXT)

# 右侧图片
add_image_placeholder(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(5.5), "目标示意图/图表")

# 第 6 页：事故统计
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "二、水务安全生产现状 - 事故统计（2021-2023）")

# 统计数据
stats = [
    "全国水利行业事故数据",
    "",
    "136 起  生产安全事故",
    "152 人  死亡",
    "28 人   受伤",
    "",
    "事故类型分布：",
    "• 水利建设项目事故 68%",
    "  高处坠落、物体打击、溺水、机械伤害",
    "",
    "• 水库及水电站运行事故 22%",
    "  设施故障引发的溃坝、触电等"
]
add_content_with_image(slide, "", stats, Inches(7), Inches(1.3), Inches(5.5), Inches(5), "事故统计图表/饼图")

# 第 7 页：安全生产特点
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "承德水务安全生产四大特点")

features = [
    ("风险点多面广", "四河之源、两库上游\n1500 条河流\n众多水库和水电站"),
    ("季节性风险突出", "冬季：冰面安全、设施冻损\n夏季：山洪、滑坡、河道行洪"),
    ("工程类型复杂", "运行类：水库、水电站\n建设类：河道治理、设施提质增效"),
    ("与生态安全关联", "京津唐重要水源地\n直接影响区域生态环境和供水安全")
]

for i, (title, text) in enumerate(features):
    left = Inches(0.6) + (i % 2) * Inches(6.8)
    top = Inches(1.3) + (i // 2) * Inches(2.8)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.3), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_LIGHT
    box.line.color.rgb = COLOR_SECONDARY
    
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_font(p, size=18, bold=True, color=COLOR_PRIMARY)
    p = tf.add_paragraph()
    p.text = text
    set_font(p, size=14, color=COLOR_TEXT)

# 第 8 页：已开展工作
content8 = [
    "一是安全生产责任体系逐步完善",
    "• 制定领导干部职责清单、科室单位职责清单、工作指导意见",
    "• 构建责任闭环，形成各负其责、齐抓共管格局",
    "",
    "二是专项整治行动深入开展",
    "• 治本攻坚三年行动、百日攻坚、冬春集中攻坚",
    "• 常态化隐患排查、双盲抽查",
    "• 事故发生率稳中有降",
    "",
    "三是安全生产基础工作不断夯实",
    "• 常态化安全警示教育和业务培训",
    "• 推进风险管控六项机制建设",
    "• 完善应急预案体系，组织应急演练"
]
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_with_image(slide, "已开展的安全生产工作", content8, Inches(7.5), Inches(1.3), Inches(5), Inches(2.5), "检查工作照片")

# 第 9 页：存在问题
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "（四）存在的主要问题")

# 监管层面
box1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(6), Inches(2.8))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(255, 240, 240)
box1.line.color.rgb = RGBColor(200, 50, 50)
tf = box1.text_frame
p = tf.paragraphs[0]
p.text = "监管层面四大短板"
set_font(p, size=18, bold=True, color=RGBColor(180, 0, 0))
p = tf.add_paragraph()
p.text = "1. 责任传导上热下冷\n2. 协同管控合力不足\n3. 执法监管刚性不够\n4. 隐患识别能力偏弱"
set_font(p, size=14, color=COLOR_TEXT)

# 执行层面
box2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.3), Inches(6), Inches(2.5))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(255, 245, 230)
box2.line.color.rgb = RGBColor(230, 150, 50)
tf = box2.text_frame
p = tf.paragraphs[0]
p.text = "执行层面三大问题"
set_font(p, size=18, bold=True, color=RGBColor(180, 80, 0))
p = tf.add_paragraph()
p.text = "1. 主体责任虚化\n2. 隐患排查整治不彻底\n3. 人员安全素养不足"
set_font(p, size=14, color=COLOR_TEXT)

add_image_placeholder(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(5.5), "问题示意图")

# 第 10 页：面临形势
content10 = [
    "（一）工作标准越来越严",
    "国家安全生产考核由党中央组织实施，考核规格、标准大幅提升",
    "",
    "（二）风险考验越来越大",
    "极端天气趋多趋频趋强趋广，降雨带北抬",
    "",
    "（三）事故发生概率增加",
    "灾后重建、国债项目、重点工程集中开工，高危作业增多",
    "",
    "（四）舆情关注度不断提高",
    "水务安全事关民生，社会关注度高，易引发负面舆情"
]
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_with_image(slide, "三、面临的形势", content10, Inches(7.5), Inches(1.3), Inches(5), Inches(5), "形势分析图")

# 第 11 页：下一步工作
content11 = [
    "一是持续健全责任体系",
    "动态调整职责清单，构建全链条、全流程、全覆盖闭环管理",
    "",
    "二是持续深化专项整治",
    "完成治本攻坚三年行动（2024-2026），重点领域拉网式排查",
    "",
    "三是持续提升防控能力",
    "推进风险管控六项机制建设，常态化培训教育",
    "",
    "四是持续完善应急保障",
    "修订应急预案，常态化应急演练，提升实战能力"
]
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_with_image(slide, "四、下一步重点工作", content11, Inches(7.5), Inches(1.3), Inches(5), Inches(5), "工作计划图/时间表")

# 第 12 页：工程施工隐患
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "五、常见安全隐患 - 工程施工环节")

hazards = [
    "1. 现场管理不规范 - 无专项方案、无证上岗",
    "2. 基坑边坡围堰风险 - 支护不到位、排水不畅",
    "3. 高处作业脚手架 - 搭设不规范、未系安全带",
    "4. 施工用电 - 电线破损、私拉乱接",
    "5. 起重机械设备 - 未检测验收、带病运行",
    "6. 消防安全 - 易燃物乱堆、动火无审批"
]

for i, h in enumerate(hazards):
    top = Inches(1.2) + i * Inches(0.95)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(7.5), Inches(0.85))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 245, 245) if i % 2 == 0 else RGBColor(255, 255, 255)
    box.line.color.rgb = RGBColor(200, 50, 50)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = h
    set_font(p, size=15, color=COLOR_TEXT)

add_image_placeholder(slide, Inches(8.5), Inches(1.2), Inches(4), Inches(5.5), "施工现场隐患照片")

# 第 13 页：设施运行隐患
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "五、常见安全隐患 - 设施运行环节")

run_hazards = [
    "1. 水库大坝 - 坝体渗漏管涌、溢洪道堵塞、监测设施异常",
    "2. 水闸泵站 - 闸门变形锈蚀、机电设备老化、违规操作",
    "3. 堤防河道 - 堤防塌陷、河道淤积、违章侵占采砂",
    "4. 供水灌区 - 管道破损渗漏、阀门失灵、水源地污染",
    "5. 管理应急 - 值班制度不落实、应急物资不足、演练流于形式"
]

for i, h in enumerate(run_hazards):
    top = Inches(1.3) + i * Inches(1)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(12), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 245, 245) if i % 2 == 0 else RGBColor(255, 255, 255)
    box.line.color.rgb = RGBColor(200, 50, 50)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = h
    set_font(p, size=16, color=COLOR_TEXT)

# 第 14 页：结束页
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_PRIMARY
bg.line.fill.background()

text = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), prs.slide_width - Inches(1.6), Inches(2))
tf = text.text_frame
p = tf.paragraphs[0]
p.text = "汇报完毕"
p.alignment = PP_ALIGN.CENTER
set_font(p, size=36, bold=True, color=RGBColor(255, 255, 255))
p = tf.add_paragraph()
p.text = "请各位领导批评指正"
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(20)
set_font(p, size=24, color=RGBColor(255, 255, 255))

# 保存
output_path = "/Users/jiyi/.openclaw/workspace/水库安全生产汇报_专业版.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print("共 14 页，带图片占位符")
