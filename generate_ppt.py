#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
水库安全生产汇报 PPT 生成脚本
政府汇报风格 - 简约正式
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 比例
prs.slide_height = Inches(7.5)

# 配色方案 - 政府汇报蓝白灰
COLOR_PRIMARY = RGBColor(25, 55, 109)    # 深蓝
COLOR_SECONDARY = RGBColor(68, 114, 196)  # 中蓝
COLOR_ACCENT = RGBColor(237, 125, 49)     # 橙色点缀
COLOR_TEXT = RGBColor(0, 0, 0)            # 黑色文字
COLOR_TEXT_LIGHT = RGBColor(80, 80, 80)   # 灰色文字
COLOR_BG = RGBColor(255, 255, 255)        # 白色背景

def set_font(paragraph, size=18, bold=False, color=COLOR_TEXT, font_name='微软雅黑'):
    """设置字体样式"""
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

def create_title_slide(prs, title, subtitle=None):
    """创建封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    
    # 顶部蓝色条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.2), prs.slide_width - Inches(1.6), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=40, bold=True, color=COLOR_PRIMARY)
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(3.8), prs.slide_width - Inches(1.6), Inches(0.8)
        )
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.alignment = PP_ALIGN.CENTER
        set_font(p, size=20, color=COLOR_TEXT_LIGHT)
    
    # 底部装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(6.8), prs.slide_width - Inches(8), Inches(0.15)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_SECONDARY
    line.line.fill.background()
    
    return slide

def create_content_slide(prs, title, content_items, subtitle=None):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.25), prs.slide_width - Inches(1.2), Inches(0.6)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    set_font(p, size=26, bold=True, color=RGBColor(255, 255, 255))
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.85), prs.slide_width - Inches(1.2), Inches(0.3)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        set_font(p, size=14, color=COLOR_SECONDARY)
    
    # 内容区域
    content_top = Inches(1.3) if not subtitle else Inches(1.5)
    content_height = prs.slide_height - content_top - Inches(0.5)
    
    content_box = slide.shapes.add_textbox(
        Inches(0.6), content_top, prs.slide_width - Inches(1.2), content_height
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.space_after = Pt(12)
        p.level = 0
        set_font(p, size=18, color=COLOR_TEXT)
    
    return slide

def create_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """创建双栏内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部标题栏
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.25), prs.slide_width - Inches(1.2), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_font(p, size=26, bold=True, color=RGBColor(255, 255, 255))
    
    # 左栏
    left_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.3), Inches(6), Inches(5.5)
    )
    tf = left_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = left_title
    set_font(p, size=20, bold=True, color=COLOR_PRIMARY)
    
    for item in left_items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.space_after = Pt(8)
        p.level = 0
        set_font(p, size=16, color=COLOR_TEXT)
    
    # 右栏
    right_box = slide.shapes.add_textbox(
        Inches(7), Inches(1.3), Inches(5.7), Inches(5.5)
    )
    tf = right_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = right_title
    set_font(p, size=20, bold=True, color=COLOR_PRIMARY)
    
    for item in right_items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.space_after = Pt(8)
        p.level = 0
        set_font(p, size=16, color=COLOR_TEXT)
    
    return slide

def create_end_slide(prs):
    """创建结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()
    
    # 文字
    text_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.5), prs.slide_width - Inches(1.6), Inches(2)
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "汇报完毕"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, size=36, bold=True, color=RGBColor(255, 255, 255))
    
    p = tf.add_paragraph()
    p.text = "请各位领导批评指正"
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(20)
    set_font(p, size=24, color=RGBColor(255, 255, 255))
    
    return slide

# ============ 开始生成 PPT ============

# 第 1 页：封面
create_title_slide(
    prs,
    "守牢安全底线 护航水务发展",
    "承德市水务局安全生产工作汇报"
)

# 第 2 页：目录
create_content_slide(
    prs,
    "汇报目录",
    [
        "一、切实加强水务安全生产工作的认识",
        "二、水务安全生产现状",
        "三、面临的形势",
        "四、下一步重点工作",
        "五、常见安全隐患"
    ]
)

# 第 3 页：重要论述核心要义
create_content_slide(
    prs,
    "一、切实加强水务安全生产工作的认识",
    [
        "书记关于安全生产重要论述核心要义：",
        "",
        "1. 树牢红线意识",
        "   发展决不能以牺牲人的生命为代价",
        "",
        "2. 健全责任体系",
        "   党政同责、一岗双责、齐抓共管、失职追责",
        "",
        "3. 压实主体责任",
        "   安全投入、管理、培训、防护、应急'五到位'",
        "",
        "4. 创新监管方式，构建长效机制",
        "5. 强化责任担当，时时放心不下"
    ]
)

# 第 4 页：法律法规体系
create_two_column_slide(
    prs,
    "核心法律法规与政策文件",
    "通用法规",
    [
        "《安全生产法》",
        "• 明确企业主体责任",
        "• 明确政府监管职责",
        "• 隐患排查治理要求",
        "• 应急处置基本原则"
    ],
    "水利专项法规",
    [
        "《水法》《防洪法》",
        "• 工程设计施工运行维护",
        "• 全生命周期安全管理",
        "• 事前防控、过程控制、事后处置",
        "",
        "省市制度",
        "• 安全生产末端落实工作机制",
        "• 安全生产违法行为行政处罚规定"
    ]
)

# 第 5 页：工作目标
create_content_slide(
    prs,
    "水务安全生产目标",
    [
        "省水利厅总体目标：",
        "坚决杜绝重特大事故",
        "遏制较大事故",
        "减少一般事故",
        "守牢安全底线",
        "",
        "我市基本底线：",
        "人员不伤亡、财产少损失",
        "",
        "确保全市水利工程建设有序推进",
        "确保各类水利设施安全稳定运行",
        "为全市经济社会发展提供坚实水安全保障"
    ]
)

# 第 6 页：事故统计数据
create_content_slide(
    prs,
    "二、水务安全生产现状 - 事故统计",
    [
        "2021-2023 年全国水利行业事故数据：",
        "",
        "共发生生产安全事故 136 起",
        "造成 152 人死亡、28 人受伤",
        "",
        "事故类型分布：",
        "• 水利建设项目事故占比 68%",
        "  主要集中在：高处坠落、物体打击、溺水、机械伤害",
        "",
        "• 水库及水电站运行事故占比 22%",
        "  主要是：设施故障引发的溃坝、触电等事故",
        "",
        "结论：水利行业安全生产风险依然突出，必须时刻绷紧安全这根弦"
    ]
)

# 第 7 页：安全生产特点
create_content_slide(
    prs,
    "承德水务安全生产四大特点",
    [
        "1. 风险点多面广",
        "   四河之源、两库上游，1500 条河流，众多水库和水电站",
        "   从山区小型水库到城区河道治理，监管覆盖范围广、难度大",
        "",
        "2. 季节性风险突出",
        "   冬季：冰面安全、设施冻损风险",
        "   夏季：山洪、滑坡、河道行洪风险",
        "",
        "3. 工程类型复杂多样",
        "   水库、水电站运行类 + 河道治理、设施提质增效建设类",
        "   老旧设施维护改造 + 新建工程施工建设",
        "",
        "4. 与生态安全紧密关联",
        "   京津唐重要水源地，直接影响区域生态环境和供水安全"
    ]
)

# 第 8 页：已开展工作
create_two_column_slide(
    prs,
    "已开展的安全生产工作",
    "责任体系建设",
    [
        "制定三项清单：",
        "• 领导干部安全生产职责清单",
        "• 机关各科室和局属单位职责清单",
        "• 安全生产工作指导意见",
        "",
        "厘清参建各方责任：",
        "• 项目法人：全过程安全负总责",
        "• 勘测设计：对成果负责",
        "• 监理单位：承担监理责任",
        "• 施工单位：全面负责本单位安全",
        "• 水行政部门：监督检查职责"
    ],
    "专项整治行动",
    [
        "持续开展专项行动：",
        "• 治本攻坚三年行动",
        "• 百日攻坚",
        "• 冬春集中攻坚",
        "• 常态化隐患排查",
        "• '双盲'抽查",
        "",
        "严厉打击：",
        "• '三违'作业",
        "• 转包、违法分包行为",
        "",
        "成效：事故发生率稳中有降"
    ]
)

# 第 9 页：存在的主要问题
create_two_column_slide(
    prs,
    "存在的主要问题",
    "监管层面四大短板",
    [
        "1. 责任传导'上热下冷'",
        "   部分县级部门、基层单位落实不到位",
        "   偏远河段、小型水库存在管控盲区",
        "",
        "2. 协同管控合力不足",
        "   市县之间、科室之间联动机制不健全",
        "",
        "3. 执法监管刚性不够",
        "   检查流于形式，多以整改代替处罚",
        "",
        "4. 隐患识别能力偏弱",
        "   基层监管人员专业素养不足"
    ],
    "执行层面三大问题",
    [
        "1. 主体责任虚化",
        "   部分单位将安全管理视为附加任务",
        "   岗位责任未落实到人",
        "",
        "2. 隐患排查整治不彻底",
        "   排查流于形式，整改不及时",
        "   存在'临时应付'心态",
        "",
        "3. 人员安全素养不足",
        "   安全意识淡薄，违规操作时有发生",
        "   特种作业人员持证上岗率未达 100%"
    ]
)

# 第 10 页：面临的形势
create_content_slide(
    prs,
    "三、面临的形势",
    [
        "1. 工作标准越来越严",
        "   国家安全生产考核由党中央组织实施",
        "   较大事故由国务院安委办挂牌督办、提级调查",
        "",
        "2. 风险考验越来越大",
        "   极端天气趋多趋频趋强趋广",
        "   降雨带'北抬'，安全生产与自然灾害风险交织",
        "",
        "3. 事故发生概率增加",
        "   灾后重建、国债项目、重点工程集中开工",
        "   高边坡、深基坑等高危作业增多",
        "   施工人员流动性大，管理难度加大",
        "",
        "4. 舆情关注度不断提高",
        "   水务安全事关民生，社会关注度高",
        "   一旦发生事故易引发负面舆情"
    ]
)

# 第 11 页：下一步重点工作
create_content_slide(
    prs,
    "四、下一步重点工作",
    [
        "1. 持续健全责任体系",
        "   动态调整职责清单，构建全链条、全流程、全覆盖闭环管理",
        "",
        "2. 持续深化专项整治",
        "   完成治本攻坚三年行动任务（2024-2026 年）",
        "   聚焦工程建设、设施运行、老旧工程改造等重点领域",
        "",
        "3. 持续提升防控能力",
        "   推进风险管控'六项机制'建设",
        "   常态化开展安全警示教育和业务培训",
        "",
        "4. 持续完善应急保障",
        "   修订完善各类突发事件应急预案",
        "   常态化开展应急演练，提升实战能力"
    ]
)

# 第 12 页：工程施工安全隐患
create_content_slide(
    prs,
    "五、常见安全隐患 - 工程施工环节",
    [
        "六大安全隐患重点排查：",
        "",
        "1. 现场管理不规范",
        "   无专项方案、不按方案施工、特种作业无证上岗",
        "",
        "2. 基坑、边坡、围堰风险",
        "   支护不到位、排水不畅、未开展位移沉降监测",
        "",
        "3. 高处作业与脚手架、模板隐患",
        "   脚手架搭设不规范、未系安全带、无防护网",
        "",
        "4. 施工用电不规范",
        "   电线破损拖地、私拉乱接、漏电保护失效",
        "",
        "5. 起重机械及设备隐患",
        "   未检测验收、限位装置失效、带病运行",
        "",
        "6. 消防安全管控不到位",
        "   易燃材料乱堆乱放、动火作业无审批、消防器材缺失"
    ]
)

# 第 13 页：设施运行安全隐患
create_content_slide(
    prs,
    "五、常见安全隐患 - 设施运行环节",
    [
        "五大安全隐患重点排查：",
        "",
        "1. 水库、大坝安全",
        "   坝体渗漏管涌、溢洪道堵塞、监测设施异常",
        "",
        "2. 水闸、泵站安全",
        "   闸门变形锈蚀、机电设备老化、违规操作",
        "",
        "3. 堤防、河道安全",
        "   堤防塌陷洞穴、河道淤积、违章侵占采砂",
        "",
        "4. 供水、灌区工程安全",
        "   管道破损渗漏、阀门失灵、水源地污染隐患",
        "",
        "5. 管理与应急保障隐患",
        "   值班巡查制度不落实、应急物资不足、演练流于形式"
    ]
)

# 第 14 页：结束页
create_end_slide(prs)

# 保存 PPT
output_path = "/Users/jiyi/.openclaw/workspace/水库安全生产汇报.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"共 14 页")
